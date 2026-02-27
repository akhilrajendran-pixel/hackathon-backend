"""
PDF / DOCX / PPTX → plain text + page info.
"""
import io
import logging
from typing import Dict, List

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_text(file_bytes: bytes, filename: str) -> Dict:
    """
    Extract text from a document.

    Returns:
        {
            "filename": str,
            "pages": [{"page": int, "text": str}, ...],
            "full_text": str,
        }
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        return _extract_pdf(file_bytes, filename)
    elif ext == "docx":
        return _extract_docx(file_bytes, filename)
    elif ext == "pptx":
        return _extract_pptx(file_bytes, filename)
    else:
        logger.warning("Unsupported file type: %s", filename)
        return {"filename": filename, "pages": [], "full_text": ""}


def _extract_pdf(data: bytes, filename: str) -> Dict:
    pages: List[Dict] = []
    doc = fitz.open(stream=data, filetype="pdf")
    for i, page in enumerate(doc):
        text = page.get_text("text")
        if text.strip():
            pages.append({"page": i + 1, "text": text})
    doc.close()
    full_text = "\n\n".join(p["text"] for p in pages)
    return {"filename": filename, "pages": pages, "full_text": full_text}


def _extract_docx(data: bytes, filename: str) -> Dict:
    doc = DocxDocument(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # DOCX doesn't have real pages; treat the whole doc as page 1
    # but split into synthetic "pages" every ~3000 chars for citation granularity
    pages: List[Dict] = []
    current_text = ""
    page_num = 1
    for para in paragraphs:
        current_text += para + "\n"
        if len(current_text) > 3000:
            pages.append({"page": page_num, "text": current_text.strip()})
            page_num += 1
            current_text = ""
    if current_text.strip():
        pages.append({"page": page_num, "text": current_text.strip()})

    full_text = "\n\n".join(p["text"] for p in pages)
    return {"filename": filename, "pages": pages, "full_text": full_text}


def _extract_pptx(data: bytes, filename: str) -> Dict:
    prs = Presentation(io.BytesIO(data))
    pages: List[Dict] = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            pages.append({"page": i + 1, "text": "\n".join(texts)})

    full_text = "\n\n".join(p["text"] for p in pages)
    return {"filename": filename, "pages": pages, "full_text": full_text}
